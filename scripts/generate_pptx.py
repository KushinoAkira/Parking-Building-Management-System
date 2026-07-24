"""
PBMS Presentation Generator
Generates a professional 16-slide PPTX for SU26SWP08 defense presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── COLOR PALETTE ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # Navy dark background
ACCENT_BLUE  = RGBColor(0x00, 0x8C, 0xFF)    # Electric blue
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)    # Cyan highlight
ACCENT_GREEN = RGBColor(0x00, 0xE0, 0x9A)    # Emerald green
ACCENT_RED   = RGBColor(0xFF, 0x4C, 0x4C)    # Live demo red
ACCENT_AMBER = RGBColor(0xFF, 0xC1, 0x07)    # Warning amber
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xBE, 0xCF)
MID_GRAY     = RGBColor(0x3A, 0x4A, 0x5C)
CARD_BG      = RGBColor(0x15, 0x2C, 0x42)    # Card background
SLIDE_W      = Inches(13.333)
SLIDE_H      = Inches(7.5)

# ─── HELPERS ──────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox

def add_bullet_slide(slide, items, left, top, width, height,
                     font_size=13, color=LIGHT_GRAY, bullet="▸ "):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return txBox

def set_slide_bg(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide, color=ACCENT_BLUE, height=0.04):
    """Thin horizontal accent bar at top"""
    add_rect(slide, 0, 0, 13.333, height, color)

def add_slide_number(slide, num, total=16):
    add_textbox(slide, f"{num:02d} / {total}", 12.3, 7.1, 1.0, 0.35,
                font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def add_section_tag(slide, tag_text, color=ACCENT_BLUE):
    rect = add_rect(slide, 0.4, 0.12, len(tag_text)*0.13+0.3, 0.28, color)
    add_textbox(slide, tag_text, 0.4, 0.12, len(tag_text)*0.13+0.3, 0.28,
                font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def add_header(slide, title, subtitle=None, title_y=0.55):
    add_textbox(slide, title, 0.4, title_y, 12.5, 0.85,
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.4, title_y + 0.75, 12.5, 0.5,
                    font_size=13, color=ACCENT_CYAN)

def make_card(slide, left, top, width, height, title, bullets,
              title_color=ACCENT_CYAN, font_size=11.5):
    add_rect(slide, left, top, width, height, CARD_BG)
    # left accent stripe
    add_rect(slide, left, top, 0.04, height, title_color)
    add_textbox(slide, title, left+0.12, top+0.1, width-0.2, 0.32,
                font_size=12, bold=True, color=title_color)
    add_bullet_slide(slide, bullets, left+0.12, top+0.42,
                     width-0.2, height-0.55, font_size=font_size,
                     color=LIGHT_GRAY, bullet="• ")

# ─── SLIDE FACTORY ────────────────────────────────────────────────────────────

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    # Gradient-like left strip
    add_rect(slide, 0, 0, 0.5, 7.5, ACCENT_BLUE)
    add_rect(slide, 0.5, 0, 0.02, 7.5, ACCENT_CYAN)
    # Top accent
    add_accent_bar(slide, ACCENT_CYAN, 0.06)
    # Logo placeholder text
    add_textbox(slide, "FPT UNIVERSITY", 1.0, 0.3, 5, 0.4,
                font_size=11, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
    add_textbox(slide, "SWP391 — Software Project", 1.0, 0.68, 8, 0.35,
                font_size=11, color=LIGHT_GRAY)
    # Main title
    add_textbox(slide, "PARKING BUILDING", 1.0, 1.5, 11.5, 1.0,
                font_size=52, bold=True, color=WHITE)
    add_textbox(slide, "MANAGEMENT SYSTEM", 1.0, 2.35, 11.5, 1.0,
                font_size=52, bold=True, color=ACCENT_CYAN)
    add_textbox(slide, "PBMS — SU26SWP08", 1.0, 3.35, 8, 0.45,
                font_size=18, color=LIGHT_GRAY)
    # Divider
    add_rect(slide, 1.0, 3.9, 8.0, 0.03, ACCENT_BLUE)
    # Team placeholder
    add_textbox(slide, "📋  Nhóm phát triển SU26SWP08  |  Tháng 7, 2026", 1.0, 4.05, 10, 0.45,
                font_size=13, color=LIGHT_GRAY)
    add_textbox(slide, "Điền tên thành viên tại đây", 1.0, 4.5, 10, 0.4,
                font_size=11, color=MID_GRAY, italic=True)

def slide_02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_BLUE)
    add_section_tag(slide, "AGENDA")
    add_header(slide, "Cấu trúc buổi bảo vệ", "Tổng thời gian: 45 phút")
    add_slide_number(slide, 2)

    blocks = [
        ("🔷 BLOCK 1", "Mở đầu & Bài toán", "~3 phút", ACCENT_BLUE,    0.5,  2.0, 5.6, 1.2),
        ("🔷 BLOCK 2", "Giải pháp & Phạm vi", "~4 phút", ACCENT_GREEN,  6.6,  2.0, 5.6, 1.2),
        ("🔷 BLOCK 3", "Kiến trúc & Kỹ thuật", "~8 phút", ACCENT_CYAN,  0.5,  3.5, 5.6, 1.2),
        ("🔴 BLOCK 4", "Live Demo", "~15 phút", ACCENT_RED,             6.6,  3.5, 5.6, 1.2),
        ("🔷 BLOCK 5", "Kết quả & Roadmap + Q&A", "~15 phút", ACCENT_AMBER, 0.5, 5.0, 11.7, 1.2),
    ]
    for tag, title, duration, color, l, t, w, h in blocks:
        add_rect(slide, l, t, w, h, CARD_BG)
        add_rect(slide, l, t, 0.06, h, color)
        add_textbox(slide, tag, l+0.18, t+0.1, w-0.3, 0.3,
                    font_size=10, bold=True, color=color)
        add_textbox(slide, title, l+0.18, t+0.35, w-0.3, 0.4,
                    font_size=15, bold=True, color=WHITE)
        add_textbox(slide, duration, l+0.18, t+0.78, w-0.3, 0.3,
                    font_size=11, color=LIGHT_GRAY)

def slide_03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_RED)
    add_section_tag(slide, "BÀI TOÁN", ACCENT_RED)
    add_header(slide, "Vấn đề của bãi đỗ xe truyền thống", "6 pain points thực tế")
    add_slide_number(slide, 3)

    problems = [
        ("❌", "Không có bản đồ slot real-time",   "Nhân viên tuần tra thủ công → tắc nghẽn giờ cao điểm"),
        ("❌", "Tính phí thủ công",                 "Tính nhẩm, thối tiền sai → thiếu minh bạch với khách"),
        ("❌", "Không có đặt chỗ trước",            "Khách VIP không có ưu tiên → trải nghiệm kém"),
        ("❌", "Quản lý vi phạm bằng sổ giấy",     "Không truy xuất lịch sử, dễ thất lạc hồ sơ"),
        ("❌", "Báo cáo doanh thu chậm trễ",        "Manager tổng hợp Excel cuối ngày, không ra quyết định kịp"),
        ("❌", "Không kiểm soát loại xe / làn",     "Xe xăng đỗ khu EV, xe tải vào khu xe máy"),
    ]
    for i, (icon, title, desc) in enumerate(problems):
        row = i // 2
        col = i % 2
        l = 0.4 + col * 6.5
        t = 2.0 + row * 1.55
        add_rect(slide, l, t, 6.3, 1.35, CARD_BG)
        add_rect(slide, l, t, 0.06, 1.35, ACCENT_RED)
        add_textbox(slide, icon + "  " + title, l+0.18, t+0.1, 6.0, 0.38,
                    font_size=13, bold=True, color=WHITE)
        add_textbox(slide, desc, l+0.18, t+0.52, 5.9, 0.7,
                    font_size=11, color=LIGHT_GRAY)

def slide_04_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_GREEN)
    add_section_tag(slide, "GIẢI PHÁP", ACCENT_GREEN)
    add_header(slide, "Tổng quan sản phẩm PBMS", "Số hóa toàn diện vòng đời một phiên đỗ xe")
    add_slide_number(slide, 4)

    # Left: scope stats
    stats = [
        ("🏢", "4 Tầng", "800+ chỗ đỗ"),
        ("🗺️", "104 Zones", "A–Z × 4 floors"),
        ("🚗", "5 Loại xe", "Motorbike, Car, EV variants"),
        ("⚡", "EV Zones E–I", "Tích hợp trạm sạc"),
    ]
    for i, (icon, val, label) in enumerate(stats):
        t = 2.0 + i * 1.2
        add_rect(slide, 0.4, t, 3.8, 1.05, CARD_BG)
        add_rect(slide, 0.4, t, 0.06, 1.05, ACCENT_GREEN)
        add_textbox(slide, icon + "  " + val, 0.6, t+0.08, 3.4, 0.38,
                    font_size=14, bold=True, color=ACCENT_GREEN)
        add_textbox(slide, label, 0.6, t+0.45, 3.4, 0.35,
                    font_size=11, color=LIGHT_GRAY)

    # Right: value chain flow
    add_textbox(slide, "Vòng đời phiên đỗ xe (End-to-End)", 4.6, 1.85, 8.4, 0.4,
                font_size=13, bold=True, color=ACCENT_CYAN)
    flow = [
        ("📅", "Đặt chỗ trước", ACCENT_BLUE),
        ("📷", "Check-in OCR", ACCENT_CYAN),
        ("🗺️", "Giám sát Real-time", ACCENT_GREEN),
        ("💰", "Tính phí & Checkout", ACCENT_AMBER),
        ("📊", "Analytics", ACCENT_BLUE),
    ]
    for i, (icon, label, clr) in enumerate(flow):
        t = 2.4 + i * 0.92
        add_rect(slide, 4.6, t, 8.3, 0.76, CARD_BG)
        add_rect(slide, 4.6, t, 0.06, 0.76, clr)
        idx_text = f"0{i+1}"
        add_textbox(slide, idx_text, 4.75, t+0.18, 0.4, 0.35,
                    font_size=14, bold=True, color=clr)
        add_textbox(slide, f"{icon}  {label}", 5.2, t+0.18, 7.5, 0.35,
                    font_size=13, bold=True, color=WHITE)
        if i < len(flow) - 1:
            add_textbox(slide, "↓", 8.5, t+0.76, 0.3, 0.3,
                        font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

def slide_05_roles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_CYAN)
    add_section_tag(slide, "USER ROLES & USE CASES")
    add_header(slide, "Đối tượng người dùng & Phân quyền RBAC", "42 Use Cases | 4 Roles | 9 Modules")
    add_slide_number(slide, 5)

    roles = [
        ("👤 ADMIN",   ACCENT_BLUE,   ["Tạo & quản lý tài khoản nội bộ", "Khóa/mở tài khoản", "Cấu hình hệ thống (grace period, VIP surcharge...)", "Xem toàn bộ audit log"]),
        ("📊 MANAGER", ACCENT_GREEN,  ["Thiết lập chính sách giá theo loại xe", "Quản lý 104 zones & 800 slots", "Xem Dashboard KPI real-time", "Xử lý incidents, quản lý gói tháng"]),
        ("🧑‍💼 STAFF",  ACCENT_CYAN,   ["Check-in xe (OCR + thủ công)", "Check-out & thu phí (Cash/EWallet)", "Xem bản đồ slot, tạo incident report", "Xác nhận / hủy đặt chỗ trước"]),
        ("🚗 DRIVER",  ACCENT_AMBER,  ["Đăng ký / Login (Email + Google OAuth)", "Đặt chỗ trước, mua vé tháng", "Nạp ví qua PayOS QR, tự check-out", "Xem lịch sử, nhận thông báo real-time"]),
    ]
    for i, (role, color, bullets) in enumerate(roles):
        col = i % 2
        row = i // 2
        l = 0.4 + col * 6.5
        t = 2.0 + row * 2.45
        add_rect(slide, l, t, 6.3, 2.25, CARD_BG)
        add_rect(slide, l, t, 6.3, 0.38, color)
        add_textbox(slide, role, l+0.15, t+0.06, 6.0, 0.3,
                    font_size=14, bold=True, color=WHITE)
        add_bullet_slide(slide, bullets, l+0.15, t+0.45, 6.0, 1.65,
                         font_size=11, color=LIGHT_GRAY)

    # UC count badges
    modules = ["Auth (5)", "Facility (6)", "Operations (6)", "Payment (4)",
               "Reservation (3)", "Subscription (3)", "Incident (3)", "Analytics (7)", "Self-service (5)"]
    add_textbox(slide, "Tổng 42 Use Cases phân theo module:", 0.4, 6.82, 12.5, 0.3,
                font_size=10, bold=True, color=ACCENT_CYAN)
    badge_text = "  |  ".join(modules)
    add_textbox(slide, badge_text, 0.4, 7.1, 12.5, 0.3,
                font_size=9.5, color=LIGHT_GRAY)

def slide_06_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_BLUE)
    add_section_tag(slide, "KIẾN TRÚC HỆ THỐNG")
    add_header(slide, "Monolithic Layered Architecture + CI/CD", "1 API | 1 DB | 3-Tier | Docker → Railway + Firebase")
    add_slide_number(slide, 6)

    # 3 tier diagram
    tiers = [
        ("CLIENT TIER",  "React 18 + Vite + Tailwind CSS\nFirebase Hosting (CDN Global)", ACCENT_CYAN,  0.4,  2.0, 3.8, 2.0),
        ("API TIER",     ".NET 10 Web API + SignalR\nEF Core 10 | Docker → Railway", ACCENT_BLUE,  4.6,  2.0, 4.0, 2.0),
        ("DATA TIER",    "PostgreSQL 14\n17 core tables | Npgsql", ACCENT_GREEN, 9.0,  2.0, 3.9, 2.0),
    ]
    arrow_positions = [(4.2, 2.75), (8.6, 2.75)]
    for (title, body, color, l, t, w, h) in tiers:
        add_rect(slide, l, t, w, h, CARD_BG)
        add_rect(slide, l, t, w, 0.38, color)
        add_textbox(slide, title, l+0.1, t+0.06, w-0.2, 0.28,
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, body, l+0.1, t+0.5, w-0.2, 1.3,
                    font_size=11.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    for (ax, ay) in arrow_positions:
        add_textbox(slide, "◀──▶", ax, ay, 0.55, 0.35,
                    font_size=13, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # External integrations row
    add_textbox(slide, "External Integrations", 0.4, 4.25, 12.5, 0.35,
                font_size=12, bold=True, color=ACCENT_CYAN)
    integrations = [
        ("🔑 Google OAuth 2.0", "Federated login cho Driver", ACCENT_BLUE),
        ("💳 PayOS Gateway",    "Payment link + HMAC Webhook", ACCENT_GREEN),
        ("📷 PaddleOCR",        "Windows local plate OCR", ACCENT_AMBER),
        ("🌐 PlateRecognizer",  "Linux/cloud OCR fallback", ACCENT_AMBER),
    ]
    for i, (name, desc, color) in enumerate(integrations):
        l = 0.4 + i * 3.25
        add_rect(slide, l, 4.65, 3.1, 0.95, CARD_BG)
        add_rect(slide, l, 4.65, 0.05, 0.95, color)
        add_textbox(slide, name, l+0.15, 4.72, 2.9, 0.3,
                    font_size=11, bold=True, color=WHITE)
        add_textbox(slide, desc, l+0.15, 5.02, 2.9, 0.45,
                    font_size=10, color=LIGHT_GRAY)

    # CI/CD
    add_textbox(slide, "CI/CD Pipeline: Push → GitHub Actions (Build + Secret Scan + Deploy) → Railway API | Firebase SPA → Smoke Test",
                0.4, 5.75, 12.5, 0.45,
                font_size=11, color=LIGHT_GRAY)

def slide_07_realtime(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_GREEN)
    add_section_tag(slide, "CORE FEATURE 1", ACCENT_GREEN)
    add_header(slide, "Real-Time Slot Map & Smart Allocation", "SignalR WebSocket | Serializable Transaction | Zero Race Condition")
    add_slide_number(slide, 7)

    # Slot status legend
    statuses = [
        ("🟢 Available",    ACCENT_GREEN),
        ("🔴 Occupied",     ACCENT_RED),
        ("🟡 Reserved",     ACCENT_AMBER),
        ("⬜ Maintenance",  MID_GRAY),
    ]
    add_textbox(slide, "Trạng thái Slot — Color Coding", 0.4, 2.0, 6.0, 0.38,
                font_size=13, bold=True, color=ACCENT_CYAN)
    for i, (label, color) in enumerate(statuses):
        add_rect(slide, 0.4 + i*1.5, 2.45, 1.35, 0.55, CARD_BG)
        add_textbox(slide, label, 0.4 + i*1.5 + 0.08, 2.52, 1.25, 0.38,
                    font_size=10.5, bold=True, color=color)

    # SignalR groups
    add_textbox(slide, "SignalR Hub — 3 Broadcast Groups", 0.4, 3.2, 6.0, 0.38,
                font_size=13, bold=True, color=ACCENT_CYAN)
    groups = [
        ('Group "all"',          'slotUpdated → cập nhật bản đồ tất cả clients'),
        ('Group "operations"',   'dashboardRefresh → KPI Manager real-time'),
        ('Group "driver:{id}"',  'Personal notification (ví, session end)'),
    ]
    for i, (grp, desc) in enumerate(groups):
        t = 3.65 + i * 0.65
        add_rect(slide, 0.4, t, 6.0, 0.55, CARD_BG)
        add_rect(slide, 0.4, t, 0.05, 0.55, ACCENT_GREEN)
        add_textbox(slide, grp, 0.6, t+0.08, 2.2, 0.35,
                    font_size=11, bold=True, color=ACCENT_GREEN)
        add_textbox(slide, desc, 2.85, t+0.08, 3.55, 0.35,
                    font_size=11, color=LIGHT_GRAY)

    # Allocation algorithm
    add_textbox(slide, "SlotAllocationService — Thuật toán cấp phát thông minh", 6.8, 2.0, 6.1, 0.4,
                font_size=13, bold=True, color=ACCENT_CYAN)
    steps = [
        ("1", "Filter",      "Chỉ xét slots đúng VehicleType"),
        ("2", "Zone rank",   "Chọn zone có tỷ lệ lấp đầy thấp nhất"),
        ("3", "VIP check",   "Slot vị trí 1 = VIP (+10.000 VND)"),
        ("4", "Lock",        "Serializable Tx — prevent concurrent assign"),
        ("5", "Broadcast",   "SignalR → slotUpdated tới group 'all'"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        t = 2.5 + i * 0.88
        add_rect(slide, 6.8, t, 6.1, 0.75, CARD_BG)
        add_rect(slide, 6.8, t, 0.4, 0.75, ACCENT_GREEN)
        add_textbox(slide, num, 6.82, t+0.18, 0.38, 0.38,
                    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, 7.28, t+0.08, 1.4, 0.3,
                    font_size=12, bold=True, color=WHITE)
        add_textbox(slide, desc, 7.28, t+0.4, 5.5, 0.28,
                    font_size=10.5, color=LIGHT_GRAY)

    # Key highlight
    add_rect(slide, 0.4, 5.9, 12.5, 0.5, RGBColor(0x05, 0x2A, 0x10))
    add_textbox(slide, "⚡  ISS-001 RESOLVED: Serializable Transaction ngăn chặn hoàn toàn race condition — 2 request đồng thời cho 1 slot, PostgreSQL serialize và reject 1 trong 2.",
                0.55, 5.94, 12.2, 0.4, font_size=10.5, color=ACCENT_GREEN)

def slide_08_checkin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_CYAN)
    add_section_tag(slide, "CORE FEATURE 2")
    add_header(slide, "Check-in / Check-out & Thanh toán số", "OCR → Auto Slot → Pricing Formula → PayOS EWallet")
    add_slide_number(slide, 8)

    # Flow: Check-in
    add_textbox(slide, "🚗  Luồng CHECK-IN", 0.4, 1.95, 6.0, 0.38,
                font_size=13, bold=True, color=ACCENT_CYAN)
    cin_steps = ["Xe đến → Staff scan OCR (PaddleOCR/PlateRec)", "Validate: không có session đang active (BR-001)",
                 "SlotAllocationService → cấp slot tối ưu", "Serializable Tx → ParkingSessions INSERT",
                 "Response: TicketCode + SlotNumber", "SignalR → slotUpdated broadcast"]
    add_bullet_slide(slide, cin_steps, 0.4, 2.38, 6.0, 2.6, font_size=11.5, color=LIGHT_GRAY)

    # Pricing formula box
    add_textbox(slide, "💰  Công thức Tính phí (PricingService)", 6.6, 1.95, 6.4, 0.38,
                font_size=13, bold=True, color=ACCENT_AMBER)
    formula_lines = [
        "hours = ceil(ExitTime - EntryTime)",
        "fee   = hours × hourly_rate",
        "if fee > daily_cap: fee = daily_cap",
        "if duration ≤ grace_period: fee = 0",
        "if active_subscription: fee = 0",
        "total = fee + vip_surcharge + penalty_fee",
    ]
    add_rect(slide, 6.6, 2.38, 6.4, 2.3, RGBColor(0x0A, 0x20, 0x30))
    for i, line in enumerate(formula_lines):
        add_textbox(slide, line, 6.75, 2.45 + i*0.35, 6.2, 0.32,
                    font_size=10.5, color=ACCENT_CYAN if i == 5 else LIGHT_GRAY)

    # Flow: Check-out
    add_textbox(slide, "🏁  Luồng CHECK-OUT", 0.4, 5.1, 6.0, 0.38,
                font_size=13, bold=True, color=ACCENT_GREEN)
    cout_steps = ["Query session (TicketCode / Plate)", "Preview phí để Staff xác nhận",
                  "Chọn PT: Cash | BankTransfer | EWallet", "EWallet: WalletService.DeductAsync() — Serializable Tx",
                  "Session → Completed | Slot → Available", "SignalR broadcast + Payment record"]
    add_bullet_slide(slide, cout_steps, 0.4, 5.52, 6.0, 1.85, font_size=11.5, color=LIGHT_GRAY)

    # Issue highlight
    add_rect(slide, 6.6, 5.1, 6.4, 1.3, RGBColor(0x1A, 0x10, 0x05))
    add_textbox(slide, "⚡  ISS-008 RESOLVED", 6.75, 5.17, 6.1, 0.3,
                font_size=11, bold=True, color=ACCENT_AMBER)
    add_textbox(slide,
                "EWallet double-deduction prevention:\nSession chỉ được Complete 1 lần (idempotency check) + Serializable Transaction bảo vệ WalletService.DeductAsync().",
                6.75, 5.5, 6.1, 0.85, font_size=10.5, color=LIGHT_GRAY)

def slide_09_payos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_AMBER)
    add_section_tag(slide, "PAYMENT & SECURITY", ACCENT_AMBER)
    add_header(slide, "Tích hợp PayOS & Bảo mật Webhook", "HMAC-SHA256 | Idempotent Webhook | Real-Time Wallet Update")
    add_slide_number(slide, 9)

    # PayOS flow
    flow = [
        ("01", "Driver", "Nhấn 'Nạp tiền', chọn mệnh giá"),
        ("02", "API",    "POST /api/portal/driver/{id}/wallet/top-up"),
        ("03", "PayOS",  "Tạo Payment Link + QR Code"),
        ("04", "Driver", "Scan QR → chuyển khoản ngân hàng"),
        ("05", "PayOS",  "POST webhook → /api/payos/webhook"),
        ("06", "API",    "Verify HMAC-SHA256 → Credit balance"),
        ("07", "SignalR","Push wallet update → driver:{id}"),
    ]
    add_textbox(slide, "Luồng Nạp tiền PayOS (7 bước)", 0.4, 1.95, 7.5, 0.38,
                font_size=13, bold=True, color=ACCENT_AMBER)
    for i, (num, actor, action) in enumerate(flow):
        t = 2.4 + i * 0.67
        clr = ACCENT_GREEN if actor == "API" else (ACCENT_AMBER if actor == "PayOS" else (ACCENT_CYAN if actor == "SignalR" else LIGHT_GRAY))
        add_rect(slide, 0.4, t, 7.5, 0.58, CARD_BG)
        add_textbox(slide, num, 0.55, t+0.13, 0.45, 0.3,
                    font_size=11, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"[{actor}]", 1.1, t+0.13, 1.3, 0.3,
                    font_size=11, bold=True, color=clr)
        add_textbox(slide, action, 2.5, t+0.13, 5.3, 0.3,
                    font_size=11, color=LIGHT_GRAY)

    # HMAC detail
    add_textbox(slide, "🔐  HMAC-SHA256 Validation (ISS-003 Resolved)", 8.2, 1.95, 4.8, 0.38,
                font_size=12, bold=True, color=ACCENT_AMBER)
    hmac_steps = [
        "Nhận webhook payload từ PayOS",
        "Sort tất cả keys alphabetically",
        "Concatenate: key=value&key=value...",
        "HMAC-SHA256 với Checksum Key",
        "So sánh với signature trong header",
        "Nếu mismatch → 401 Unauthorized",
    ]
    add_rect(slide, 8.2, 2.4, 4.8, 3.5, RGBColor(0x15, 0x10, 0x05))
    for i, step in enumerate(hmac_steps):
        add_textbox(slide, f"{i+1}. {step}", 8.35, 2.5 + i*0.53, 4.6, 0.45,
                    font_size=11, color=LIGHT_GRAY)

    # Root cause callout
    add_rect(slide, 8.2, 6.02, 4.8, 0.72, RGBColor(0x2A, 0x18, 0x00))
    add_textbox(slide, "Root Cause ISS-003:", 8.35, 6.08, 4.5, 0.28,
                font_size=10, bold=True, color=ACCENT_AMBER)
    add_textbox(slide, "Payload serialize sai thứ tự key → custom flattener sort alphabetical trước khi hash.",
                8.35, 6.35, 4.5, 0.35, font_size=10, color=LIGHT_GRAY)

def slide_10_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_BLUE)
    add_section_tag(slide, "SECURITY & QUALITY")
    add_header(slide, "Bảo mật & Chất lượng Engineering", "5 lớp bảo vệ | NFRs đạt chuẩn | 10 Issues Resolved")
    add_slide_number(slide, 10)

    layers = [
        ("🔑", "JWT + BCrypt",          "JWT 24h | BCrypt cost=11 | OAuth2 Driver only",    ACCENT_BLUE),
        ("🚦", "Rate Limiting",          "8 req/phút/IP cho login endpoint",                   ACCENT_BLUE),
        ("🔒", "Serializable Tx",        "Check-in & EWallet deduct: IsolationLevel.Serializable", ACCENT_RED),
        ("📋", "HMAC Webhook",           "PayOS signature verification alphabetical-sort",    ACCENT_AMBER),
        ("🔍", "Gitleaks Secret Scan",   "Mọi PR: Full git history scan, block nếu phát hiện secret", ACCENT_GREEN),
    ]
    add_textbox(slide, "5 Lớp Bảo mật", 0.4, 1.9, 6.2, 0.38,
                font_size=13, bold=True, color=ACCENT_CYAN)
    for i, (icon, title, desc, color) in enumerate(layers):
        t = 2.35 + i * 0.9
        add_rect(slide, 0.4, t, 6.2, 0.78, CARD_BG)
        add_rect(slide, 0.4, t, 0.05, 0.78, color)
        add_textbox(slide, f"{icon}  {title}", 0.6, t+0.08, 5.8, 0.3,
                    font_size=12, bold=True, color=color)
        add_textbox(slide, desc, 0.6, t+0.42, 5.8, 0.3,
                    font_size=10.5, color=LIGHT_GRAY)

    # NFR + Issues
    add_textbox(slide, "Non-Functional Requirements (Đạt)", 6.8, 1.9, 6.1, 0.38,
                font_size=13, bold=True, color=ACCENT_CYAN)
    nfrs = [
        ("⚡", "API Response", "< 500 ms (read)"),
        ("⏱️", "Check-in/out", "< 2 giây end-to-end"),
        ("👥", "Concurrent Users", "50 active users"),
        ("🛡️", "Financial Integrity", "Serializable isolation"),
        ("🤖", "Background Workers", "Reservation & Subscription expiry"),
    ]
    for i, (icon, label, val) in enumerate(nfrs):
        t = 2.35 + i * 0.73
        add_rect(slide, 6.8, t, 6.1, 0.62, CARD_BG)
        add_textbox(slide, f"{icon}  {label}", 7.0, t+0.08, 3.5, 0.28,
                    font_size=12, bold=True, color=WHITE)
        add_textbox(slide, val, 10.5, t+0.08, 2.3, 0.28,
                    font_size=12, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.RIGHT)

    add_rect(slide, 6.8, 6.05, 6.1, 0.5, RGBColor(0x05, 0x1A, 0x08))
    add_textbox(slide, "✅  10 Issues resolved — Critical: ISS-001 (Race condition), ISS-008 (Double deduction) | High: ISS-002 (OCR Linux), ISS-003 (HMAC), ISS-004 (OAuth prod)",
                6.95, 6.1, 5.9, 0.4, font_size=10, color=ACCENT_GREEN)

def slide_11_demo_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_RED)
    add_section_tag(slide, "LIVE DEMO — 15 PHÚT", ACCENT_RED)
    add_header(slide, "Kịch bản Demo hệ thống", "Hệ thống đã được deploy trên Railway + Firebase (Production)")
    add_slide_number(slide, 11)

    demos = [
        ("01", "3 phút",  "Manager Dashboard & Slot Map Real-time",  "KPI live, bản đồ 800 chỗ màu sắc, SignalR push",       ACCENT_BLUE),
        ("02", "4 phút",  "Staff Check-in xe",                       "OCR nhận diện biển số → auto slot → TicketCode",        ACCENT_CYAN),
        ("03", "3 phút",  "Driver: Nạp ví & Đặt chỗ trước",        "PayOS QR → webhook → balance | VIP Reservation",        ACCENT_AMBER),
        ("04", "4 phút",  "Staff Check-out & Tính phí tự động",     "Preview fee → EWallet deduct → slot freed → SignalR",    ACCENT_GREEN),
        ("05", "1 phút",  "Analytics Dashboard",                     "Revenue trend 30 ngày, Zone occupancy report",          ACCENT_BLUE),
    ]
    for i, (num, duration, title, desc, color) in enumerate(demos):
        t = 2.0 + i * 1.02
        add_rect(slide, 0.4, t, 12.5, 0.88, CARD_BG)
        add_rect(slide, 0.4, t, 0.55, 0.88, color)
        add_textbox(slide, num, 0.42, t+0.22, 0.55, 0.38,
                    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, duration, 1.1, t+0.05, 1.1, 0.3,
                    font_size=10, bold=True, color=color)
        add_textbox(slide, title, 1.1, t+0.32, 5.5, 0.3,
                    font_size=13, bold=True, color=WHITE)
        add_textbox(slide, desc, 6.8, t+0.22, 6.0, 0.42,
                    font_size=11, color=LIGHT_GRAY)

def make_demo_backdrop(prs, slide_num, title, caption, total=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x06, 0x10, 0x18))
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT_RED)
    # LIVE badge
    add_rect(slide, 0.3, 0.2, 1.5, 0.55, ACCENT_RED)
    add_textbox(slide, "🔴  LIVE DEMO", 0.35, 0.27, 1.45, 0.38,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, 2.1, 0.25, 10.5, 0.55,
                font_size=22, bold=True, color=WHITE)
    add_textbox(slide, caption, 0.3, 0.95, 12.7, 0.4,
                font_size=12, color=ACCENT_CYAN)
    add_slide_number(slide, slide_num, total)
    # Large "Demo in progress" watermark area
    add_textbox(slide, "[ Màn hình demo trực tiếp trên trình duyệt ]",
                0, 3.2, 13.333, 1.0,
                font_size=22, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)
    return slide

def slide_12_demo_manager(prs):
    make_demo_backdrop(prs, 12,
        "Manager Dashboard & Real-Time Slot Map",
        "Giám sát 800 chỗ đỗ theo thời gian thực | SignalR WebSocket | KPI Dashboard")

def slide_13_demo_checkin(prs):
    make_demo_backdrop(prs, 13,
        "Staff Check-in & OCR Nhận diện Biển số",
        "PaddleOCR → Plate text → Auto-allocate slot → Serializable Transaction → TicketCode")

def slide_14_demo_driver(prs):
    make_demo_backdrop(prs, 14,
        "Driver: Nạp Ví PayOS & Đặt chỗ trước (VIP Reservation)",
        "PayOS QR Code → Bank transfer → HMAC Webhook → Credit balance → SignalR push wallet update")

def slide_15_demo_checkout(prs):
    make_demo_backdrop(prs, 15,
        "Staff Check-out & Tính phí Tự động",
        "Fee preview: ceil(hours)×rate | Daily cap | Grace period | EWallet deduct → Session Complete → Slot freed")

def slide_16_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, ACCENT_GREEN)
    add_section_tag(slide, "KẾT QUẢ & ROADMAP", ACCENT_GREEN)
    add_header(slide, "Kết quả đạt được & Hướng phát triển", "Summer 2026 — SU26SWP08")
    add_slide_number(slide, 16)

    # Achievements
    add_textbox(slide, "✅  Thành tựu dự án", 0.4, 1.9, 6.0, 0.38,
                font_size=14, bold=True, color=ACCENT_GREEN)
    achievements = [
        "42 Use Cases | 9 Modules | 4 Roles (RBAC)",
        "17 DB tables | 800+ slots | 104 zones",
        "3 external integrations (PayOS, Google OAuth, OCR dual-engine)",
        "10 issues resolved — bao gồm 2 CRITICAL",
        "API < 500ms | Check-in/out < 2s | 50 concurrent users",
        "CI/CD + DevSecOps (Gitleaks) hoàn chỉnh",
        "Production deploy: Railway API + Firebase Hosting",
    ]
    add_bullet_slide(slide, achievements, 0.4, 2.35, 6.0, 3.5, font_size=12, color=LIGHT_GRAY)

    # Roadmap
    add_textbox(slide, "🚀  Hướng phát triển tương lai", 6.8, 1.9, 6.1, 0.38,
                font_size=14, bold=True, color=ACCENT_CYAN)
    roadmap = [
        ("P1", "📱 Mobile App (React Native)", "Driver UX tốt hơn, native QR scan", ACCENT_BLUE),
        ("P1", "📹 CCTV Auto Check-in",        "Auto nhận diện xe vào cổng — zero staff", ACCENT_GREEN),
        ("P2", "🤖 AI Slot Recommendation",     "Personalized zone theo lịch sử Driver", ACCENT_CYAN),
        ("P2", "⏰ Report Scheduler (CRON)",    "Tự động tạo daily snapshot", ACCENT_AMBER),
        ("P3", "🌐 Multi-facility (SaaS)",      "Mở rộng chuỗi bãi đỗ xe", ACCENT_BLUE),
        ("P3", "⚡ EV Charging Integration",    "Quản lý sạc Zones E–I", ACCENT_GREEN),
    ]
    for i, (priority, title, desc, color) in enumerate(roadmap):
        t = 2.35 + i * 0.82
        add_rect(slide, 6.8, t, 6.1, 0.7, CARD_BG)
        add_rect(slide, 6.8, t, 0.5, 0.7, color)
        add_textbox(slide, priority, 6.82, t+0.18, 0.5, 0.3,
                    font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, 7.38, t+0.06, 5.4, 0.28,
                    font_size=12, bold=True, color=WHITE)
        add_textbox(slide, desc, 7.38, t+0.38, 5.4, 0.25,
                    font_size=10.5, color=LIGHT_GRAY)

    # Thank you footer
    add_rect(slide, 0, 7.0, 13.333, 0.5, CARD_BG)
    add_textbox(slide, "Cảm ơn Hội đồng đã lắng nghe!  |  Nhóm SU26SWP08 sẵn sàng trả lời câu hỏi & demo live.",
                0.4, 7.08, 12.5, 0.35, font_size=12, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# ─── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Generating PBMS presentation...")
    steps = [
        (slide_01_cover,       "Slide 01 — Cover"),
        (slide_02_agenda,      "Slide 02 — Agenda"),
        (slide_03_problem,     "Slide 03 — Problem Statement"),
        (slide_04_overview,    "Slide 04 — Product Overview"),
        (slide_05_roles,       "Slide 05 — User Roles & Use Cases"),
        (slide_06_architecture,"Slide 06 — Architecture & Tech Stack"),
        (slide_07_realtime,    "Slide 07 — Real-Time Slot Map"),
        (slide_08_checkin,     "Slide 08 — Check-in/out & Payment"),
        (slide_09_payos,       "Slide 09 — PayOS & Webhook Security"),
        (slide_10_security,    "Slide 10 — Security & Quality"),
        (slide_11_demo_plan,   "Slide 11 — Demo Plan"),
        (slide_12_demo_manager,"Slide 12 — [DEMO] Manager Dashboard"),
        (slide_13_demo_checkin,"Slide 13 — [DEMO] Check-in OCR"),
        (slide_14_demo_driver, "Slide 14 — [DEMO] Driver Wallet"),
        (slide_15_demo_checkout,"Slide 15 — [DEMO] Checkout"),
        (slide_16_results,     "Slide 16 — Results & Roadmap"),
    ]
    for fn, label in steps:
        fn(prs)
        print(f"  [OK] {label}")

    out_path = r"f:\FPT_material\2026\Summer\SWP391\Project\Parking-Building-Management-System\docs\PBMS_SU26SWP08_Presentation.pptx"
    prs.save(out_path)
    print(f"\n[DONE] Saved: {out_path}")
    print(f"       16 slides | 45-minute format | Navy dark theme")


if __name__ == "__main__":
    main()
